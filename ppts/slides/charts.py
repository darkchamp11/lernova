"""
charts.py - Native PPTX charts.
"""

from pptx.chart.data import CategoryChartData, ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.util import Inches, Pt
from .tokens import *
from .helpers import add_text, add_eyebrow_header

def add_bar_chart(slide, x, y, width, height, title, categories, series_data):
    """
    Add a grouped bar chart.
    categories: list of str (e.g. ['2021', '2022', '2023'])
    series_data: dict of series_name -> list of values (e.g. {'Series A': [10, 20, 30]})
    """
    chart_data = CategoryChartData()
    chart_data.categories = categories
    
    for name, values in series_data.items():
        chart_data.add_series(name, values)
        
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, width, height, chart_data
    ).chart

    # Styling the chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.font.name = SANS
    chart.legend.font.size = Pt(10)
    chart.legend.font.color.rgb = OLIVE
    chart.legend.include_in_layout = False
    
    chart.has_title = False # We'll use a slide header instead

    # Axis styling
    category_axis = chart.category_axis
    category_axis.has_major_gridlines = False
    category_axis.tick_labels.font.name = SANS
    category_axis.tick_labels.font.size = Pt(10)
    category_axis.tick_labels.font.color.rgb = CHARCOAL
    
    value_axis = chart.value_axis
    value_axis.has_major_gridlines = True
    value_axis.major_gridlines.format.line.color.rgb = BORDER_CREAM
    value_axis.tick_labels.font.name = SANS
    value_axis.tick_labels.font.size = Pt(10)
    value_axis.tick_labels.font.color.rgb = STONE

    # Color the series
    base_colors = [SERIES_1, SERIES_2, SERIES_3, SERIES_4, SERIES_5, SERIES_6]
    for idx, series in enumerate(chart.series):
        fill = series.format.fill
        fill.solid()
        fill.fore_color.rgb = base_colors[idx % len(base_colors)]


def add_donut_chart(slide, x, y, width, height, title, categories, values):
    """
    Add a donut chart using base colors.
    categories: list of str
    values: list of float/int
    """
    chart_data = ChartData()
    chart_data.categories = categories
    chart_data.add_series(title, values)
    
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT, x, y, width, height, chart_data
    ).chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.font.name = SANS
    chart.legend.font.size = Pt(10)
    chart.legend.font.color.rgb = OLIVE
    
    chart.has_title = False
    
    # Color the slices
    base_colors = [SERIES_1, SERIES_2, SERIES_3, SERIES_4, SERIES_5, SERIES_6]
    # In doughnut charts, we color the individual points
    for idx, point in enumerate(chart.series[0].points):
        fill = point.format.fill
        fill.solid()
        fill.fore_color.rgb = base_colors[idx % len(base_colors)]
        
    return chart
