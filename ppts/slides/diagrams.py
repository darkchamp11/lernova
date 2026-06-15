"""
diagrams.py - Native PPTX shapes representing diagrams (architecture, flowchart, etc.)
"""

from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from .tokens import *
from .helpers import add_card, add_text, add_line, add_brand_left_bar

def add_architecture_diagram(slide, left_in, top_in):
    """
    Draws a native PPTX architecture diagram.
    """
    # Legend
    leg_y = Inches(top_in + 3.2)
    add_text(slide, "LEGEND", Inches(left_in), leg_y, Inches(1), Inches(0.2), font=MONO, size=8, color=STONE)
    add_line(slide, Inches(left_in), leg_y + Inches(0.25), Inches(10), color=BORDER_WARM, weight_pt=1)
    
    # Legend items
    add_card(slide, Inches(left_in), leg_y + Inches(0.4), Inches(0.15), Inches(0.1), fill=BRAND_TINT, border=BRAND, border_weight=1)
    add_text(slide, "Focal · Origin", Inches(left_in + 0.25), leg_y + Inches(0.35), Inches(1.5), Inches(0.2), font=SANS, size=9, color=OLIVE)
    
    add_card(slide, Inches(left_in + 1.8), leg_y + Inches(0.4), Inches(0.15), Inches(0.1), fill=IVORY, border=NEAR_BLACK, border_weight=1)
    add_text(slide, "Backend", Inches(left_in + 2.05), leg_y + Inches(0.35), Inches(1.5), Inches(0.2), font=SANS, size=9, color=OLIVE)
    
    # Drawing Nodes
    # 1. External (User)
    n1_x = left_in + 0.5
    n1_y = top_in + 1.5
    add_card(slide, Inches(n1_x), Inches(n1_y), Inches(1.5), Inches(0.8), fill=PARCHMENT, border=STONE, border_weight=1)
    add_text(slide, "USER", Inches(n1_x), Inches(n1_y + 0.1), Inches(1.5), Inches(0.2), font=MONO, size=7, color=STONE, align=PP_ALIGN.CENTER)
    add_text(slide, "Browser Client", Inches(n1_x), Inches(n1_y + 0.3), Inches(1.5), Inches(0.25), font=SANS, size=11, bold=True, color=NEAR_BLACK, align=PP_ALIGN.CENTER)
    
    # 2. CDN/Edge
    n2_x = n1_x + 2.0
    n2_y = n1_y
    add_card(slide, Inches(n2_x), Inches(n2_y), Inches(1.5), Inches(0.8), fill=PARCHMENT, border=STONE, border_weight=1)
    add_text(slide, "EDGE", Inches(n2_x), Inches(n2_y + 0.1), Inches(1.5), Inches(0.2), font=MONO, size=7, color=STONE, align=PP_ALIGN.CENTER)
    add_text(slide, "Leaflet Maps", Inches(n2_x), Inches(n2_y + 0.3), Inches(1.5), Inches(0.25), font=SANS, size=11, bold=True, color=NEAR_BLACK, align=PP_ALIGN.CENTER)
    
    # 3. Origin (FOCAL)
    n3_x = n2_x + 2.0
    n3_y = n1_y
    add_card(slide, Inches(n3_x), Inches(n3_y), Inches(1.7), Inches(0.8), fill=BRAND_TINT, border=BRAND, border_weight=1.5)
    add_text(slide, "ORIGIN", Inches(n3_x), Inches(n3_y + 0.1), Inches(1.7), Inches(0.2), font=MONO, size=7, color=BRAND, align=PP_ALIGN.CENTER)
    add_text(slide, "FastAPI Server", Inches(n3_x), Inches(n3_y + 0.3), Inches(1.7), Inches(0.25), font=SANS, size=11, bold=True, color=NEAR_BLACK, align=PP_ALIGN.CENTER)
    add_text(slide, "SSE · Streams", Inches(n3_x), Inches(n3_y + 0.55), Inches(1.7), Inches(0.2), font=MONO, size=8, color=OLIVE, align=PP_ALIGN.CENTER)
    
    # 4. Storage / ML Models
    n4_x = n3_x + 2.2
    n4_y = n1_y - 0.6
    add_card(slide, Inches(n4_x), Inches(n4_y), Inches(1.5), Inches(0.8), fill=IVORY, border=NEAR_BLACK, border_weight=1)
    add_text(slide, "INFERENCE", Inches(n4_x), Inches(n4_y + 0.1), Inches(1.5), Inches(0.2), font=MONO, size=7, color=NEAR_BLACK, align=PP_ALIGN.CENTER)
    add_text(slide, "LSTM Models", Inches(n4_x), Inches(n4_y + 0.3), Inches(1.5), Inches(0.25), font=SANS, size=11, bold=True, color=NEAR_BLACK, align=PP_ALIGN.CENTER)
    
    n5_x = n3_x + 2.2
    n5_y = n1_y + 0.6
    add_card(slide, Inches(n5_x), Inches(n5_y), Inches(1.5), Inches(0.8), fill=IVORY, border=NEAR_BLACK, border_weight=1)
    add_text(slide, "STORE", Inches(n5_x), Inches(n5_y + 0.1), Inches(1.5), Inches(0.2), font=MONO, size=7, color=NEAR_BLACK, align=PP_ALIGN.CENTER)
    add_text(slide, "InMemoryStore", Inches(n5_x), Inches(n5_y + 0.3), Inches(1.5), Inches(0.25), font=SANS, size=11, bold=True, color=NEAR_BLACK, align=PP_ALIGN.CENTER)

    # Connections
    add_line(slide, Inches(n1_x + 1.5), Inches(n1_y + 0.4), Inches(0.5), color=OLIVE, weight_pt=1.5)
    add_line(slide, Inches(n2_x + 1.5), Inches(n2_y + 0.4), Inches(0.5), color=BRAND, weight_pt=2.0)
    
    add_line(slide, Inches(n3_x + 1.7), Inches(n1_y + 0.2), Inches(0.5), color=OLIVE, weight_pt=1.5)
    add_line(slide, Inches(n3_x + 1.7), Inches(n1_y + 0.6), Inches(0.5), color=OLIVE, weight_pt=1.5)


def add_flowchart_diagram(slide, left_in, top_in):
    """
    Draws a native PPTX flowchart for the Risk Engine.
    """
    # 1. Start Pill
    s_x, s_y = left_in + 4.0, top_in + 0.2
    add_card(slide, Inches(s_x), Inches(s_y), Inches(1.5), Inches(0.4), fill=PARCHMENT, border=STONE, border_weight=1)
    add_text(slide, "Forecast Snapshot", Inches(s_x), Inches(s_y + 0.1), Inches(1.5), Inches(0.25), font=SANS, size=11, bold=True, color=NEAR_BLACK, align=PP_ALIGN.CENTER)
    
    # Arrow down
    add_line(slide, Inches(s_x + 0.75), Inches(s_y + 0.4), Inches(0.5), color=OLIVE, weight_pt=1.5) # Will rotate later or just use vertical connector
    conn = slide.shapes.add_connector(1, Inches(s_x + 0.75), Inches(s_y + 0.4), Inches(s_x + 0.75), Inches(s_y + 0.9))
    conn.line.color.rgb = OLIVE
    
    # 2. Check Step
    c_x, c_y = s_x - 0.25, s_y + 0.9
    add_card(slide, Inches(c_x), Inches(c_y), Inches(2.0), Inches(0.8), fill=IVORY, border=NEAR_BLACK, border_weight=1)
    add_text(slide, "EVALUATE RISK", Inches(c_x), Inches(c_y + 0.1), Inches(2.0), Inches(0.2), font=MONO, size=7, color=STONE, align=PP_ALIGN.CENTER)
    add_text(slide, "Calculate Ratio", Inches(c_x), Inches(c_y + 0.3), Inches(2.0), Inches(0.25), font=SANS, size=11, bold=True, color=NEAR_BLACK, align=PP_ALIGN.CENTER)
    
    conn2 = slide.shapes.add_connector(1, Inches(c_x + 1.0), Inches(c_y + 0.8), Inches(c_x + 1.0), Inches(c_y + 1.3))
    conn2.line.color.rgb = OLIVE
    
    # 3. Decision (Diamond substitute - just a card for now to avoid complex geometry)
    d_x, d_y = c_x, c_y + 1.3
    add_card(slide, Inches(d_x), Inches(d_y), Inches(2.0), Inches(0.8), fill=BRAND_TINT, border=BRAND, border_weight=1.5)
    add_text(slide, "BRANCH", Inches(d_x), Inches(d_y + 0.1), Inches(2.0), Inches(0.2), font=MONO, size=7, color=BRAND, align=PP_ALIGN.CENTER)
    add_text(slide, "Ratio >= 1.65?", Inches(d_x), Inches(d_y + 0.3), Inches(2.0), Inches(0.25), font=SANS, size=11, bold=True, color=NEAR_BLACK, align=PP_ALIGN.CENTER)
    
    # Branch 1 (Yes)
    conn3 = slide.shapes.add_connector(1, Inches(d_x), Inches(d_y + 0.4), Inches(d_x - 1.0), Inches(d_y + 0.4))
    conn3.line.color.rgb = BRAND
    conn3.line.width = Pt(1.5)
    
    y_x, y_y = d_x - 2.5, d_y + 0.1
    add_card(slide, Inches(y_x), Inches(y_y), Inches(1.5), Inches(0.6), fill=IVORY, border=ERROR_RED, border_weight=1.5)
    add_text(slide, "HIGH ALERT", Inches(y_x), Inches(y_y + 0.1), Inches(1.5), Inches(0.25), font=SANS, size=11, bold=True, color=ERROR_RED, align=PP_ALIGN.CENTER)
    
    # Branch 2 (No)
    conn4 = slide.shapes.add_connector(1, Inches(d_x + 2.0), Inches(d_y + 0.4), Inches(d_x + 3.0), Inches(d_y + 0.4))
    conn4.line.color.rgb = OLIVE
    
    n_x, n_y = d_x + 3.0, d_y + 0.1
    add_card(slide, Inches(n_x), Inches(n_y), Inches(1.5), Inches(0.6), fill=IVORY, border=NEAR_BLACK, border_weight=1)
    add_text(slide, "NORMAL / WATCH", Inches(n_x), Inches(n_y + 0.1), Inches(1.5), Inches(0.25), font=SANS, size=11, bold=True, color=NEAR_BLACK, align=PP_ALIGN.CENTER)

