"""
helpers.py - Core drawing primitives for PPTX slides.

Low-level functions: blank_slide, add_text, add_line, add_card, add_tag,
add_bullet_list, brand_left_bar, etc. Every template in templates.py
is built from these primitives.
"""

from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from .tokens import *


# ─────────────────────────────────────────────────────────
# Slide foundation
# ─────────────────────────────────────────────────────────

def blank_slide(prs, bg_color=PARCHMENT):
    """Create a blank slide with solid background fill."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return slide


# ─────────────────────────────────────────────────────────
# Text primitives
# ─────────────────────────────────────────────────────────

def add_text(slide, text, left, top, width, height,
             font=SANS, size=18, bold=False, italic=False,
             color=NEAR_BLACK, align=PP_ALIGN.LEFT,
             vanchor=MSO_ANCHOR.TOP, line_spacing=None):
    """Add a single-paragraph textbox."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = vanchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = Pt(line_spacing)
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_multiline_text(slide, lines, left, top, width, height,
                       font=SANS, size=14, color=DARK_WARM,
                       align=PP_ALIGN.LEFT, line_spacing=18):
    """Add a textbox with multiple paragraphs, each line as a paragraph."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.line_spacing = Pt(line_spacing)
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return tb


# ─────────────────────────────────────────────────────────
# Decorative elements
# ─────────────────────────────────────────────────────────

def add_line(slide, left, top, width, color=BRAND, weight_pt=1):
    """Add a horizontal connector line."""
    line = slide.shapes.add_connector(1, left, top, left + width, top)
    line.line.color.rgb = color
    line.line.width = Pt(weight_pt)
    return line


def add_brand_left_bar(slide, left, top, height, width=Pt(3)):
    """Signature: 2.5pt brand-colored vertical bar (section-title accent)."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = BRAND
    bar.line.fill.background()
    bar.shadow.inherit = False
    bar.rotation = 0
    return bar


def add_dot(slide, cx, cy, radius=Inches(0.04), color=BRAND):
    """Small circular dot — used in eyebrow headers."""
    d = radius * 2
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - radius, cy - radius, d, d)
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    dot.shadow.inherit = False
    return dot


# ─────────────────────────────────────────────────────────
# Cards and containers
# ─────────────────────────────────────────────────────────

def add_card(slide, left, top, width, height,
             fill=IVORY, border=BORDER_CREAM, border_weight=0.5, radius=None):
    """Ivory card with cream border."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = border
    card.line.width = Pt(border_weight)
    card.shadow.inherit = False
    return card


def add_featured_card(slide, left, top, width, height):
    """Featured card with brand border (highlighted container)."""
    return add_card(slide, left, top, width, height,
                    fill=BRAND_TINT, border=BRAND, border_weight=1)


def add_dark_card(slide, left, top, width, height):
    """Dark card for contrast sections."""
    return add_card(slide, left, top, width, height,
                    fill=DARK_SURFACE, border=DARK_SURFACE, border_weight=0)


# ─────────────────────────────────────────────────────────
# Tag badges
# ─────────────────────────────────────────────────────────

def add_tag(slide, text, left, top, width=None, height=Inches(0.3),
            bg=TAG_LIGHTEST, text_color=BRAND):
    """Tag badge — solid hex background, uppercase label."""
    if width is None:
        width = Inches(max(0.8, len(text) * 0.11 + 0.3))
    tag_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       left, top, width, height)
    tag_shape.fill.solid()
    tag_shape.fill.fore_color.rgb = bg
    tag_shape.line.fill.background()
    tag_shape.shadow.inherit = False

    add_text(slide, text.upper(), left, top, width, height,
             font=SANS, size=9, bold=True, color=text_color,
             align=PP_ALIGN.CENTER, vanchor=MSO_ANCHOR.MIDDLE)
    return tag_shape


# ─────────────────────────────────────────────────────────
# Bullet lists (em-dash style, editorial)
# ─────────────────────────────────────────────────────────

def add_bullet_list(slide, bullets, left, top, width, item_height=0.55,
                    font_size=14, dash_color=BRAND, text_color=CHARCOAL):
    """Em-dash bullet list — editorial style."""
    for i, bullet in enumerate(bullets):
        y = top + Inches(i * item_height)
        add_text(slide, "—", left, y, Inches(0.3), Inches(0.4),
                 font=SANS, size=font_size, color=dash_color)
        add_text(slide, bullet, left + Inches(0.35), y, width - Inches(0.35), Inches(0.4),
                 font=SANS, size=font_size, color=text_color)


def add_numbered_list(slide, items, left, top, width, item_height=0.55,
                      font_size=14, num_color=BRAND, text_color=CHARCOAL):
    """Numbered list with brand-colored indices."""
    for i, item in enumerate(items):
        y = top + Inches(i * item_height)
        add_text(slide, f"{i+1:02d}", left, y, Inches(0.4), Inches(0.4),
                 font=SERIF, size=font_size + 2, color=num_color)
        add_text(slide, item, left + Inches(0.45), y, width - Inches(0.45), Inches(0.4),
                 font=SANS, size=font_size, color=text_color)


# ─────────────────────────────────────────────────────────
# Eyebrow header (dot + label + rule)
# ─────────────────────────────────────────────────────────

def add_eyebrow_header(slide, eyebrow, title, y_start=Inches(0.5)):
    """Section header: dot + uppercase eyebrow + title + rule."""
    # Dot
    add_dot(slide, Inches(1.1), y_start + Inches(0.08))
    # Eyebrow text
    add_text(slide, eyebrow.upper(),
             Inches(1.25), y_start - Inches(0.05), Inches(10), Inches(0.3),
             font=SANS, size=11, color=STONE, bold=True)
    # Horizontal rule
    add_line(slide, Inches(1.0), y_start + Inches(0.35), Inches(11.33),
             color=BORDER_WARM, weight_pt=0.8)
    # Title
    add_text(slide, title,
             Inches(1.0), y_start + Inches(0.55), CONTENT_W, Inches(0.9),
             font=SERIF, size=34, color=NEAR_BLACK)
    return y_start + Inches(1.6)


# ─────────────────────────────────────────────────────────
# Page number footer
# ─────────────────────────────────────────────────────────

def add_page_number(slide, num, total=None):
    """Bottom-right page number in stone color."""
    label = f"{num:02d}" if total is None else f"{num:02d} / {total:02d}"
    add_text(slide, label,
             Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.3),
             font=SANS, size=10, color=STONE, align=PP_ALIGN.RIGHT)
    # Deck mark
    add_text(slide, "FLOOD FORECASTING",
             Inches(0.8), Inches(7.0), Inches(3), Inches(0.3),
             font=SANS, size=8, color=STONE, bold=True)
